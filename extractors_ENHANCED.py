"""
ARGO - Document Extractors and Chunking
CRITICAL: This is the ONLY place where text chunking happens
"""
import os
from pathlib import Path
from typing import List, Dict, Any, Optional
import hashlib

from langchain.text_splitter import RecursiveCharacterTextSplitter

from core.config import get_config
from core.logger import get_logger

logger = get_logger("extractors")


class ChunkingStrategy:
    """
    Intelligent chunking strategy selector
    Determines optimal chunk size based on document characteristics
    """
    
    @staticmethod
    def select_chunk_size(text_length: int, file_type: str) -> Dict[str, int]:
        """
        Select optimal chunk size based on document size and type
        
        Args:
            text_length: Total length of text in characters
            file_type: Type of file (pdf, docx, txt, etc.)
        
        Returns:
            Dict with chunk_size and overlap
        """
        config = get_config()
        
        # Get configuration
        min_size = config.get("rag.chunking.min_chunk_size", 500)
        default_size = config.get("rag.chunking.default_chunk_size", 1000)
        max_size = config.get("rag.chunking.max_chunk_size", 1500)
        overlap_ratio = config.get("rag.chunking.chunk_overlap_ratio", 0.2)
        
        # Adjust based on document size
        if text_length < 5000:
            chunk_size = min_size
        elif text_length < 20000:
            chunk_size = default_size
        else:
            chunk_size = max_size
        
        # Adjust based on file type
        if file_type in ["xlsx", "csv"]:
            # Structured data: smaller chunks
            chunk_size = int(chunk_size * 0.7)
        elif file_type in ["code", "py", "js"]:
            # Code: preserve structure
            chunk_size = int(chunk_size * 0.8)
        
        overlap = int(chunk_size * overlap_ratio)
        
        return {
            "chunk_size": chunk_size,
            "chunk_overlap": overlap
        }


def extract_and_chunk(
    file_path: str,
    file_type: str,
    metadata: Optional[Dict[str, Any]] = None
) -> List[Dict[str, Any]]:
    """
    CENTRALIZED extraction and chunking function
    
    This is the ONLY function that should create RecursiveCharacterTextSplitter.
    All other modules MUST use this function for chunking.
    
    Args:
        file_path: Path to file
        file_type: Type of file (pdf, docx, txt, xlsx, etc.)
        metadata: Additional metadata to attach to chunks
    
    Returns:
        List of chunk dicts with text and metadata
    """
    logger.debug(f"Extracting and chunking file: {file_path} ({file_type})")
    
    # 1. Extract raw text
    text = extract_raw_text(file_path, file_type)
    
    if not text or not text.strip():
        logger.warning(f"No text extracted from {file_path}")
        return []
    
    # 2. Select chunking strategy
    strategy = ChunkingStrategy.select_chunk_size(len(text), file_type)
    
    logger.debug(
        f"Chunking strategy: size={strategy['chunk_size']}, "
        f"overlap={strategy['chunk_overlap']}"
    )
    
    # 3. Get hierarchical separators
    config = get_config()
    separators = config.get(
        "rag.chunking.separators",
        ["\n\n", "\n", ". ", " ", ""]
    )
    
    # 4. Create splitter (ONLY place this happens!)
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=strategy["chunk_size"],
        chunk_overlap=strategy["chunk_overlap"],
        separators=separators,
        length_function=len
    )
    
    # 5. Split text into chunks
    chunks = splitter.split_text(text)
    
    # 6. Create chunk metadata
    base_metadata = metadata or {}
    file_hash = _compute_file_hash(file_path)
    
    chunk_dicts = []
    for i, chunk_text in enumerate(chunks):
        chunk_dict = {
            "content": chunk_text,  # Changed from "text" to "content" for UI compatibility
            "text": chunk_text,      # Keep "text" for backward compatibility
            "metadata": {
                **base_metadata,
                "chunk_index": i,
                "total_chunks": len(chunks),
                "chunk_size": len(chunk_text),
                "file_hash": file_hash,
                "file_type": file_type,
                "extraction_method": "extract_and_chunk"
            }
        }
        chunk_dicts.append(chunk_dict)
    
    logger.info(
        f"Created {len(chunk_dicts)} chunks from {Path(file_path).name} "
        f"({len(text)} chars)"
    )
    
    return chunk_dicts


def extract_raw_text(file_path: str, file_type: str) -> str:
    """
    Extract raw text from file based on type
    
    Args:
        file_path: Path to file
        file_type: Type of file
    
    Returns:
        Extracted text
    """
    file_type = file_type.lower().strip('.')
    
    extractors = {
        'txt': _extract_txt,
        'md': _extract_txt,
        'pdf': _extract_pdf,
        'doc': _extract_doc,
        'docx': _extract_docx,
        'xlsx': _extract_xlsx,
        'xls': _extract_xlsx,
        'csv': _extract_csv,
        'ppt': _extract_ppt,
        'pptx': _extract_ppt,
        'mpp': _extract_mpp,
        'xer': _extract_xer,
        'png': _extract_image,
        'jpg': _extract_image,
        'jpeg': _extract_image,
        'gif': _extract_image,
        'bmp': _extract_image,
        'tiff': _extract_image,
    }
    
    extractor = extractors.get(file_type, _extract_txt)
    
    try:
        text = extractor(file_path)
        return text or ""
    except Exception as e:
        logger.error(f"Error extracting text from {file_path}: {e}")
        return ""


def _extract_txt(file_path: str) -> str:
    """Extract from plain text file"""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()


def _extract_pdf(file_path: str) -> str:
    """Extract from PDF"""
    try:
        import PyPDF2
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text
    except ImportError:
        logger.warning("PyPDF2 not installed, cannot extract PDF")
        return ""


def _extract_docx(file_path: str) -> str:
    """Extract from DOCX"""
    try:
        from docx import Document
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    except ImportError:
        logger.warning("python-docx not installed, cannot extract DOCX")
        return ""


def _extract_xlsx(file_path: str) -> str:
    """Extract from Excel"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text = ""
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text += f"\n--- Sheet: {sheet_name} ---\n"
            
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text += row_text + "\n"
        
        return text
    except ImportError:
        logger.warning("openpyxl not installed, cannot extract XLSX")
        return ""


def _extract_csv(file_path: str) -> str:
    """Extract from CSV"""
    try:
        import csv
        text = ""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            for row in reader:
                text += " | ".join(row) + "\n"
        return text
    except Exception as e:
        logger.error(f"Error reading CSV: {e}")
        return ""


def _extract_doc(file_path: str) -> str:
    """Extract from DOC (legacy Word format)"""
    try:
        import subprocess
        import tempfile
        # Try antiword if available
        result = subprocess.run(
            ['antiword', file_path],
            capture_output=True,
            text=True,
            timeout=30
        )
        if result.returncode == 0:
            return result.stdout
        else:
            logger.warning(f"antiword failed, trying textract for {file_path}")
            # Fallback to textract
            import textract
            text = textract.process(file_path).decode('utf-8', errors='ignore')
            return text
    except (ImportError, FileNotFoundError, subprocess.TimeoutExpired):
        logger.warning("DOC extraction tools not available (antiword or textract needed)")
        return f"[DOC file: {Path(file_path).name} - extraction not available]"
    except Exception as e:
        logger.error(f"Error extracting DOC: {e}")
        return f"[DOC file: {Path(file_path).name}]"


def _extract_ppt(file_path: str) -> str:
    """Extract from PPT/PPTX (PowerPoint)"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = ""

        for i, slide in enumerate(prs.slides):
            text += f"\n--- Slide {i+1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return text
    except ImportError:
        logger.warning("python-pptx not installed, cannot extract PPT/PPTX")
        return f"[PowerPoint file: {Path(file_path).name} - python-pptx needed]"
    except Exception as e:
        logger.error(f"Error extracting PPT: {e}")
        return f"[PowerPoint file: {Path(file_path).name}]"


def _extract_mpp(file_path: str) -> str:
    """Extract from MPP (Microsoft Project)"""
    try:
        # MPP is complex binary format, try mpxj if available
        import subprocess
        import tempfile
        import json

        # Try using mpxj (Java library) via command line
        temp_json = tempfile.mktemp(suffix='.json')
        result = subprocess.run(
            ['mpxjconvert', file_path, temp_json],
            capture_output=True,
            timeout=60
        )

        if result.returncode == 0 and os.path.exists(temp_json):
            with open(temp_json, 'r') as f:
                data = json.load(f)
            os.remove(temp_json)

            # Extract relevant text from project data
            text = f"Project: {data.get('project_name', 'Unknown')}\n"
            text += f"Tasks: {len(data.get('tasks', []))}\n"
            for task in data.get('tasks', [])[:100]:  # Limit to first 100 tasks
                text += f"- {task.get('name', '')}\n"
            return text
        else:
            logger.warning(f"mpxjconvert not available or failed for {file_path}")
            return f"[MS Project file: {Path(file_path).name} - mpxj tool needed for extraction]"
    except Exception as e:
        logger.error(f"Error extracting MPP: {e}")
        return f"[MS Project file: {Path(file_path).name}]"


def _extract_xer(file_path: str) -> str:
    """Extract from XER (Primavera P6)"""
    try:
        # XER is text-based format, can be parsed directly
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()

        # Extract key information
        text = f"XER Project File: {Path(file_path).name}\n\n"

        # Parse basic structure
        lines = content.split('\n')
        current_table = None

        for line in lines:
            line = line.strip()
            if line.startswith('%T'):
                # Table definition
                current_table = line.split('\t')[1] if len(line.split('\t')) > 1 else 'Unknown'
                text += f"\n--- {current_table} ---\n"
            elif line.startswith('%R'):
                # Record data
                if current_table in ['TASK', 'PROJECT', 'PROJWBS', 'TASKRSRC']:
                    fields = line.split('\t')[1:]
                    text += " | ".join(fields[:10]) + "\n"  # First 10 fields

        return text
    except Exception as e:
        logger.error(f"Error extracting XER: {e}")
        return f"[Primavera XER file: {Path(file_path).name}]"


def _extract_image(file_path: str) -> str:
    """Extract text from images using OCR"""
    try:
        from PIL import Image
        import pytesseract

        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)

        if not text.strip():
            return f"[Image file: {Path(file_path).name} - no text detected]"

        return f"[Image: {Path(file_path).name}]\n{text}"
    except ImportError:
        logger.warning("PIL or pytesseract not installed, cannot extract text from images")
        return f"[Image file: {Path(file_path).name} - OCR not available]"
    except Exception as e:
        logger.error(f"Error extracting image: {e}")
        return f"[Image file: {Path(file_path).name}]"


def _compute_file_hash(file_path: str) -> str:
    """Compute SHA256 hash of file for change detection"""
    sha256 = hashlib.sha256()
    
    try:
        with open(file_path, 'rb') as f:
            for chunk in iter(lambda: f.read(8192), b""):
                sha256.update(chunk)
        return sha256.hexdigest()
    except Exception as e:
        logger.error(f"Error computing hash for {file_path}: {e}")
        return ""


def get_file_info(file_path: str) -> Dict[str, Any]:
    """
    Get file information without extracting content
    
    Returns:
        Dict with size, type, hash, etc.
    """
    path = Path(file_path)
    
    if not path.exists():
        return {}
    
    stat = path.stat()
    
    return {
        "filename": path.name,
        "extension": path.suffix.lower().strip('.'),
        "size_bytes": stat.st_size,
        "size_kb": stat.st_size / 1024,
        "size_mb": stat.st_size / (1024 * 1024),
        "modified_timestamp": stat.st_mtime,
        "file_hash": _compute_file_hash(file_path)
    }


# Rule enforcement
def _validate_no_direct_splitter_usage():
    """
    This function serves as documentation:
    
    RULE: No other module should create RecursiveCharacterTextSplitter directly.
    ALL chunking must go through extract_and_chunk().
    
    Benefits:
    1. Consistent chunk sizes across project
    2. Centralized strategy tuning
    3. Uniform metadata
    4. Easy to debug and optimize
    5. Single source of truth
    """
    pass
